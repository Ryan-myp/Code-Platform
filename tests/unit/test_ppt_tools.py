"""v15 PPT 增强单测：4 类模板库 + 段落级结构化渲染 + 主题色板统一。

覆盖：
- PPT_TEMPLATES：4 类模板结构完整（色板六色/字体/结构原则）
- _build_ppt_system_prompt：按模板注入专属结构原则与段落级输出要求
- _build_pptx_file：段落级结构（level/emphasis）渲染、模板色板生效、字体统一、空大纲兜底
- _parse_ppt_outline：保留段落级 dict 条目
"""

import json
import sys
from pathlib import Path

BACKEND = str(Path(__file__).resolve().parents[2] / "backend")
if BACKEND not in sys.path:
    sys.path.insert(0, BACKEND)


def _sample_outline():
    return {
        "meta": {"storyline": "从问题到方案", "visual_theme": "商务蓝", "estimated_duration": "10"},
        "slides": [
            {"type": "cover", "title": "AI 商业化之路", "subtitle": "2026 战略汇报"},
            {
                "type": "content",
                "title": "市场增长强劲",
                "content": [
                    {"text": "AI 市场三年翻三倍", "level": 0, "emphasis": "strong"},
                    {"text": "2025 年规模达 1200 亿", "level": 1, "emphasis": "normal"},
                    {"text": "数据来源：行业白皮书", "level": 1, "emphasis": "quote"},
                ],
            },
            {"type": "toc", "title": "目录", "content": ["背景", "方案", "数据", "下一步"]},
        ],
    }


class TestTemplates:
    def test_four_templates_with_palette(self):
        import extended_api

        assert set(extended_api.PPT_TEMPLATES) == {"business", "roadshow", "teaching", "marketing"}
        for tid, tpl in extended_api.PPT_TEMPLATES.items():
            assert tpl["name"]
            assert tpl["desc"]
            assert tpl["font"]
            assert tpl["principles"]
            pal = tpl["palette"]
            assert set(pal) == {"dark", "accent", "accent_light", "gray", "text", "white"}
            for rgb in pal.values():
                assert len(rgb) == 3
                assert all(0 <= c <= 255 for c in rgb)

    def test_templates_palettes_differ(self):
        import extended_api

        accents = {tid: tpl["palette"]["accent"] for tid, tpl in extended_api.PPT_TEMPLATES.items()}
        assert len(set(accents.values())) == 4, "四类模板主色应互不相同"


class TestSystemPrompt:
    def test_prompt_injects_template_principles(self):
        import extended_api

        roadshow = extended_api._build_ppt_system_prompt("roadshow")
        assert "融资路演" in roadshow
        assert "TAM/SAM/SOM" in roadshow
        assert "LTV/CAC" in roadshow

        teaching = extended_api._build_ppt_system_prompt("teaching")
        assert "教学课件" in teaching
        assert "练习/互动" in teaching

        marketing = extended_api._build_ppt_system_prompt("marketing")
        assert "营销方案" in marketing
        assert "KPI" in marketing

        business = extended_api._build_ppt_system_prompt("business")
        assert "商务汇报" in business
        assert "黄金结构" in business

    def test_prompt_requires_paragraph_structure(self):
        import extended_api

        prompt = extended_api._build_ppt_system_prompt("business")
        assert "level 0" in prompt
        assert "emphasis" in prompt
        assert '"level": 0' in prompt  # JSON schema 段落级示例

    def test_unknown_template_falls_back_to_business(self):
        import extended_api

        prompt = extended_api._build_ppt_system_prompt("not_exist")
        assert "商务汇报" in prompt


class TestBuildPptx:
    def test_paragraph_structure_rendered(self, monkeypatch, tmp_path):
        import extended_api
        from pptx import Presentation
        from pptx.dml.color import RGBColor

        monkeypatch.setattr(extended_api, "PPTX_DIR", str(tmp_path))
        path = extended_api._build_pptx_file("AI 商业化之路", _sample_outline(), "business")

        prs = Presentation(path)
        assert len(prs.slides) == 3

        # 内容页：主论点加粗（strong）且为强调色；引用条目斜体
        slide2 = prs.slides[1]
        all_text = []
        for sh in slide2.shapes:
            if sh.has_text_frame:
                all_text.append(sh.text_frame.text)
        joined = "\n".join(all_text)
        assert "市场增长强劲" in joined
        assert "AI 市场三年翻三倍" in joined
        assert "2025 年规模达 1200 亿" in joined
        assert "行业白皮书" in joined

        strong_run = None
        quote_run = None
        for sh in slide2.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if "AI 市场三年翻三倍" in run.text:
                        strong_run = run
                    if "行业白皮书" in run.text:
                        quote_run = run
        assert strong_run is not None and strong_run.font.bold is True
        assert strong_run.font.color.rgb == RGBColor(0x4F, 0x46, 0xE5)  # business accent
        assert quote_run is not None and quote_run.font.italic is True

    def test_template_palette_applied_to_cover(self, monkeypatch, tmp_path):
        import extended_api
        from pptx import Presentation
        from pptx.dml.color import RGBColor

        monkeypatch.setattr(extended_api, "PPTX_DIR", str(tmp_path))
        biz = extended_api._build_pptx_file("T", _sample_outline(), "business")
        road = extended_api._build_pptx_file("T", _sample_outline(), "roadshow")

        biz_cover_bg = Presentation(biz).slides[0].shapes[0].fill.fore_color.rgb
        road_cover_bg = Presentation(road).slides[0].shapes[0].fill.fore_color.rgb
        assert biz_cover_bg == RGBColor(0x1B, 0x26, 0x3B)  # business dark
        assert road_cover_bg == RGBColor(0x16, 0x0E, 0x2B)  # roadshow dark
        assert biz_cover_bg != road_cover_bg

    def test_font_unified_across_runs(self, monkeypatch, tmp_path):
        import extended_api
        from pptx import Presentation

        monkeypatch.setattr(extended_api, "PPTX_DIR", str(tmp_path))
        path = extended_api._build_pptx_file("T", _sample_outline(), "teaching")

        prs = Presentation(path)
        fonts = set()
        for slide in prs.slides:
            for sh in slide.shapes:
                if not sh.has_text_frame:
                    continue
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        fonts.add(run.font.name)
        assert fonts == {"Microsoft YaHei"}

    def test_empty_outline_fallback(self, monkeypatch, tmp_path):
        import extended_api
        from pptx import Presentation

        monkeypatch.setattr(extended_api, "PPTX_DIR", str(tmp_path))
        path = extended_api._build_pptx_file("空演示", {"meta": {}, "slides": []}, "business")
        prs = Presentation(path)
        assert len(prs.slides) == 1
        text = " ".join(sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame)
        assert "空演示" in text

    def test_legacy_string_content_still_works(self, monkeypatch, tmp_path):
        import extended_api
        from pptx import Presentation

        outline = {"meta": {}, "slides": [{"type": "content", "title": "旧格式", "content": ["要点1", "要点2"]}]}
        monkeypatch.setattr(extended_api, "PPTX_DIR", str(tmp_path))
        path = extended_api._build_pptx_file("T", outline, "business")
        text = " ".join(
            sh.text_frame.text for sh in Presentation(path).slides[0].shapes if sh.has_text_frame
        )
        assert "要点1" in text and "要点2" in text


class TestParseOutline:
    def test_keeps_paragraph_dict_items(self):
        import extended_api

        raw = json.dumps(_sample_outline(), ensure_ascii=False)
        data = extended_api._parse_ppt_outline(raw)
        content = data["slides"][1]["content"]
        assert isinstance(content[0], dict)
        assert content[0]["level"] == 0
        assert content[0]["emphasis"] == "strong"

    def test_invalid_input_returns_empty(self):
        import extended_api

        assert extended_api._parse_ppt_outline("not json") == {"meta": {}, "slides": []}
        assert extended_api._parse_ppt_outline("") == {"meta": {}, "slides": []}
