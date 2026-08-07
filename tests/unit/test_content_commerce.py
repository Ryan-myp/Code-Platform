"""内容创作商业化升级验证。

覆盖：
- 文案/翻译/PPT 三工厂：POST 创建异步任务（task_id）→ handler 执行 → 记录落库（带 user_id）
- 历史接口用户隔离：普通用户仅见自己的记录，admin 全量；删除接口归属校验
- PPT 商业化：大纲 JSON 容错解析 → 真实 PPTX 文件生成（可读回/页数匹配）→ 下载端点
- 存量库迁移：老表缺 user_id/file_path 列时自动补列
"""

import asyncio
import json
import os
from unittest.mock import patch

import pytest
from fastapi import HTTPException

import extended_api
from task_queue import _handlers, create_task, get_task


def _user(uid: str, name: str, role: str = "user") -> dict:
    return {"user_id": uid, "username": name, "role": role}


def test_content_handlers_registered(setup_test_db):
    """三个工厂的异步任务处理器必须已注册（含用户并发限制）。"""
    missing = {"copywriting_generate", "translation_translate", "ppt_generate"} - set(_handlers)
    assert not missing, f"未注册的任务类型: {missing}"


def test_copywriting_async_flow(setup_test_db, claim_and_run):
    """文案：异步任务执行后落库（带用户归属），任务状态 success。"""
    with patch("extended_api.call_llm", return_value="# 测试文案\n\n卖点总结"):
        task = create_task(
            "copywriting_generate",
            {"type": "marketing", "title": "新品上市", "prompt": "写个文案", "user_id": "u-alice"},
            username="alice",
            user_id="u-alice",
            role="user",
        )
        assert task["status"] == "pending"
        claim_and_run(task["id"])
        assert get_task(task["id"])["status"] == "success"

    from common.db import get_db

    conn = get_db()
    try:
        rows = conn.execute("SELECT * FROM copywriting_tasks").fetchall()
        assert len(rows) == 1
        assert rows[0]["user_id"] == "u-alice"
        assert rows[0]["result"] == "# 测试文案\n\n卖点总结"
    finally:
        conn.close()


def test_copywriting_post_endpoint_creates_task(setup_test_db):
    """POST 接口直接返回 task_id（不再同步阻塞返回 result）。"""

    async def _call():
        return await extended_api.generate_copywriting(
            extended_api.CopywritingRequest(type="marketing", title="t", prompt="需求"),
            _user("u1", "alice"),
        )

    resp = asyncio.run(_call())
    assert resp["ok"] is True
    assert resp["task_id"].startswith("task_")
    assert resp["status"] == "pending"
    # 空 prompt 校验
    with pytest.raises(HTTPException) as exc:
        asyncio.run(
            extended_api.generate_copywriting(
                extended_api.CopywritingRequest(type="marketing", prompt="  "),
                _user("u1", "alice"),
            )
        )
    assert exc.value.status_code == 400


def test_translation_async_flow(setup_test_db, claim_and_run):
    """翻译：异步任务执行后落库（带用户归属）。"""
    with patch("extended_api.call_llm", return_value="Hello world"):
        task = create_task(
            "translation_translate",
            {"source_lang": "中文", "target_lang": "English", "text": "你好世界", "user_id": "u-bob"},
            username="bob",
            user_id="u-bob",
            role="user",
        )
        claim_and_run(task["id"])
        assert get_task(task["id"])["status"] == "success"

    from common.db import get_db

    conn = get_db()
    try:
        rows = conn.execute("SELECT * FROM translations").fetchall()
        assert len(rows) == 1
        assert rows[0]["user_id"] == "u-bob"
        assert rows[0]["result"] == "Hello world"
    finally:
        conn.close()


_PPT_OUTLINE = json.dumps(
    {
        "meta": {"storyline": "从问题到方案", "visual_theme": "商务蓝", "estimated_duration": "10"},
        "slides": [
            {
                "type": "cover",
                "title": "AI 商业化之路",
                "subtitle": "2026 年度汇报",
                "content": [],
                "notes": "开场问好",
            },
            {"type": "toc", "title": "目录", "content": ["背景", "方案", "数据"], "notes": ""},
            {
                "type": "content",
                "title": "核心结论",
                "content": ["要点一", "要点二"],
                "chart_suggestion": "柱状图",
                "notes": "强调数据",
            },
            {"type": "thanks", "title": "谢谢", "subtitle": "Q&A", "content": [], "notes": ""},
        ],
    },
    ensure_ascii=False,
)


def test_ppt_async_flow_with_pptx_file(setup_test_db, claim_and_run, monkeypatch, tmp_path):
    """PPT：异步任务 → 大纲解析 → 真实 PPTX 文件生成 → 记录含下载路径。"""
    monkeypatch.setattr(extended_api, "PPTX_DIR", str(tmp_path))
    with patch("extended_api.call_llm", return_value=_PPT_OUTLINE):
        task = create_task(
            "ppt_generate",
            {"title": "AI 商业化之路", "outline": "参考大纲", "user_id": "u-carl"},
            username="carl",
            user_id="u-carl",
            role="user",
        )
        claim_and_run(task["id"])
        assert get_task(task["id"])["status"] == "success"

    from common.db import get_db

    conn = get_db()
    try:
        rows = conn.execute("SELECT * FROM ppt_generations").fetchall()
        assert len(rows) == 1
        row = rows[0]
        assert row["user_id"] == "u-carl"
        assert row["file_path"].startswith("/api/ppt/download/ppt_")
        slides = json.loads(row["slides"])
        assert len(slides["slides"]) == 4
    finally:
        conn.close()

    # PPTX 文件真实可读回，页数匹配
    from pptx import Presentation

    fname = row["file_path"].rsplit("/", 1)[-1]
    pptx_path = os.path.join(str(tmp_path), fname)
    assert os.path.exists(pptx_path)
    prs = Presentation(pptx_path)
    assert len(prs.slides) == 4
    # 首页为封面版式（深色背景存在）
    first_text = " ".join(sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame)
    assert "AI 商业化之路" in first_text

    # 下载端点：返回 FileResponse 且文件存在
    resp = asyncio.run(extended_api.download_ppt(fname, _user("u-carl", "carl")))
    assert os.path.exists(resp.path)
    assert "presentationml.presentation" in resp.media_type


@pytest.mark.parametrize(
    "raw,expect_slides",
    [
        # 正常 JSON
        ('{"meta":{},"slides":[{"title":"a"}]}', 1),
        # ```json 代码块包裹
        ('```json\n{"meta":{},"slides":[{"title":"a"},{"title":"b"}]}\n```', 2),
        # 前后杂音文本
        ('好的，以下是生成结果：\n{"meta":{},"slides":[]}\n以上供参考', 0),
        # 完全非 JSON（LLM 跑飞）
        ("抱歉，我无法生成。", 0),
        # 空字符串
        ("", 0),
    ],
)
def test_ppt_outline_parse_tolerant(raw, expect_slides):
    """大纲解析必须容错：代码块包裹/杂音文本/非 JSON 均不抛异常。"""
    data = extended_api._parse_ppt_outline(raw)
    assert isinstance(data["slides"], list)
    assert len(data["slides"]) == expect_slides


def test_history_isolation_and_delete_scope(setup_test_db, claim_and_run):
    """历史接口用户隔离：普通用户仅见自己的记录；admin 全量；删除归属校验。"""
    with patch("extended_api.call_llm", return_value="A 的文案"):
        ta = create_task(
            "copywriting_generate",
            {"type": "marketing", "title": "a", "prompt": "p", "user_id": "u-a"},
            username="alice",
            user_id="u-a",
            role="user",
        )
        claim_and_run(ta["id"])
    with patch("extended_api.call_llm", return_value="B 的文案"):
        tb = create_task(
            "copywriting_generate",
            {"type": "marketing", "title": "b", "prompt": "p", "user_id": "u-b"},
            username="bob",
            user_id="u-b",
            role="user",
        )
        claim_and_run(tb["id"])

    # alice 仅见自己的 1 条
    alice_rows = asyncio.run(extended_api.list_copywriting_history(_user("u-a", "alice")))
    assert len(alice_rows) == 1 and alice_rows[0]["title"] == "a"

    # admin 全量可见
    admin_rows = asyncio.run(extended_api.list_copywriting_history(_user("admin", "admin", role="admin")))
    assert len(admin_rows) == 2

    # alice 尝试删除 bob 的记录 → 不生效（归属校验）
    bob_task_id = asyncio.run(extended_api.list_copywriting_history(_user("u-b", "bob")))[0]["id"]
    asyncio.run(extended_api.delete_copywriting(bob_task_id, _user("u-a", "alice")))
    remaining = asyncio.run(extended_api.list_copywriting_history(_user("admin", "admin", role="admin")))
    assert len(remaining) == 2

    # alice 删除自己的记录 → 生效
    alice_task_id = asyncio.run(extended_api.list_copywriting_history(_user("u-a", "alice")))[0]["id"]
    asyncio.run(extended_api.delete_copywriting(alice_task_id, _user("u-a", "alice")))
    remaining = asyncio.run(extended_api.list_copywriting_history(_user("admin", "admin", role="admin")))
    assert len(remaining) == 1


def test_legacy_table_columns_migrated(setup_test_db):
    """存量库兼容：老表缺 user_id/file_path 列时，迁移函数自动补列（幂等）。"""
    from common.db import get_db

    conn = get_db()
    try:
        # 模拟老库：删除 user_id 列（SQLite 3.35+ 支持 DROP COLUMN）
        for table in ("copywriting_tasks", "translations", "ppt_generations", "excel_operations"):
            conn.execute(f"ALTER TABLE {table} DROP COLUMN user_id")
        conn.execute("ALTER TABLE ppt_generations DROP COLUMN file_path")
        conn.commit()
    finally:
        conn.close()

    extended_api._ensure_content_user_columns()

    conn = get_db()
    try:
        for table in ("copywriting_tasks", "translations", "ppt_generations", "excel_operations"):
            cols = [r[1] for r in conn.execute(f"PRAGMA table_info({table})").fetchall()]
            assert "user_id" in cols, f"{table} 缺 user_id 列"
        cols = [r[1] for r in conn.execute("PRAGMA table_info(ppt_generations)").fetchall()]
        assert "file_path" in cols
    finally:
        conn.close()

    # 幂等：重复执行不报错
    extended_api._ensure_content_user_columns()


def test_copywriting_user_limit(setup_test_db):
    """文案工厂用户并发限制：同用户最多 2 个活跃任务。"""
    from fastapi import HTTPException

    for _ in range(2):
        create_task(
            "copywriting_generate",
            {"type": "marketing", "title": "x", "prompt": "p", "user_id": "u-l"},
            username="lim",
            user_id="u-l",
        )
    with pytest.raises(HTTPException) as exc:
        create_task(
            "copywriting_generate",
            {"type": "marketing", "title": "x", "prompt": "p", "user_id": "u-l"},
            username="lim",
            user_id="u-l",
        )
    assert exc.value.status_code == 429
